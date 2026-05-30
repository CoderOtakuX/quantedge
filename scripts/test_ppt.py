"""Test PPT generation end-to-end."""
import requests
import json
import sys

payload = {
    "verdict": "BUY",
    "confidence": 75,
    "summary": "Strong fundamentals with solid revenue growth.",
    "bull_case": "Market leader in energy sector with strong diversification",
    "bear_case": "High debt levels and regulatory risks",
    "outlook_1_2yr": "Positive outlook driven by Jio and retail growth",
    "key_risks": ["Crude oil price volatility", "Telecom competition"],
    "alternatives": ["ONGC", "BPCL"],
    "sentiment": "Positive",
    "sentiment_score": 72,
    "catalysts": ["5G rollout", "Retail expansion"],
    "news_summary": "Overall positive sentiment driven by business diversification"
}

ticker = sys.argv[1] if len(sys.argv) > 1 else "RELIANCE"
print(f"Testing PPT generation for {ticker}...")

r = requests.post(f"http://localhost:8000/api/report/{ticker}", json=payload, timeout=60)
print(f"Status: {r.status_code}")
print(f"Content-Type: {r.headers.get('Content-Type')}")
print(f"Content-Length: {len(r.content)} bytes")

if r.status_code == 200:
    fname = f"test_{ticker}.pptx"
    with open(fname, "wb") as f:
        f.write(r.content)
    print(f"Saved {fname} successfully!")
    
    # Verify it's a valid PPTX
    from pptx import Presentation
    prs = Presentation(fname)
    print(f"Slides count: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        print(f"  Slide {i+1}: {texts[:3]}...")
else:
    print(f"ERROR: {r.text}")
