import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation

p = Presentation(r'C:\Users\Admin\Downloads\IDEA_QuantEdge_Report.pptx')

for i, slide in enumerate(p.slides):
    print(f"\n=== Slide {i+1} ({len(slide.shapes)} shapes) ===")
    for s in slide.shapes:
        if s.has_text_frame:
            text = s.text_frame.text[:80].replace('\n', ' ')
            print(f"  Text: {text}")
        elif s.has_table:
            table = s.table
            print(f"  Table: {len(table.rows)} rows x {len(table.columns)} cols")
            for r in range(min(3, len(table.rows))):
                row_data = [table.cell(r, c).text[:15] for c in range(len(table.columns))]
                print(f"    Row {r}: {row_data}")
        elif s.has_chart:
            print(f"  Chart: {s.chart.chart_type}")
        else:
            print(f"  Shape: {s.shape_type}")
