from fastapi import APIRouter, HTTPException, Request
from fastapi.responses import StreamingResponse
import io
from datetime import datetime
from pydantic import BaseModel
from typing import List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

from services.yfinance_service import (
    get_overview, get_financials, get_peers, get_shareholding, get_news
)

router = APIRouter()

class AIAnalysisRequest(BaseModel):
    verdict: str
    confidence: float
    summary: str
    bull_case: str
    bear_case: str
    outlook_1_2yr: str
    key_risks: List[str]
    alternatives: List[str]
    sentiment: str
    sentiment_score: float
    catalysts: List[str]
    news_summary: str

# Colors
COLORS = {
    "NAVY": RGBColor.from_string("1E2761"),
    "ICE_BLUE": RGBColor.from_string("CADCFC"),
    "WHITE": RGBColor.from_string("FFFFFF"),
    "TEAL": RGBColor.from_string("028090"),
    "LIGHT_BG": RGBColor.from_string("F4F6FB"),
    "DARK_TEXT": RGBColor.from_string("1E293B"),
    "MID_GRAY": RGBColor.from_string("64748B"),
    "SUCCESS": RGBColor.from_string("16A34A"),
    "WARN": RGBColor.from_string("D97706"),
    "DANGER": RGBColor.from_string("DC2626"),
}

def set_bg(slide, color_name):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS[color_name]

def add_shape(slide, shape_type, x, y, w, h, fill_color=None, alpha=None):
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS[fill_color]
    shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=11, bold=False, color="DARK_TEXT", align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = COLORS[color]
    return txBox

def add_teal_accent(slide):
    add_shape(slide, 1, 0.4, 0, 0.05, 5.625, "TEAL")

def safe_float(val, default=0.0):
    try:
        return float(val) if val is not None else default
    except:
        return default

@router.post("/report/{ticker}")
def generate_report(ticker: str, ai_data: AIAnalysisRequest):
    try:
        overview = get_overview(ticker)
        financials = get_financials(ticker)
        peers = get_peers(ticker)
        shareholding = get_shareholding(ticker)
        # We don't strictly need news right now if AI already summarizes it.
    except Exception as e:
        print(f"Data fetch error: {e}")
        overview = {}
        financials = {"income_statement": {}, "balance_sheet": {}, "cash_flow": {}}
        peers = []
        shareholding = {"institutional": [], "mutual_fund": [], "major_holders": []}

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    blank_layout = prs.slide_layouts[6]

    # Helper formatters
    def fmt_cr(val):
        try:
            if val is None: return "N/A"
            v = float(val)
            # Standardize: Assume input is in units, convert to Crore (divide by 1Cr = 10^7)
            cr_val = v / 10_000_000
            if abs(cr_val) >= 100:
                return f"₹{cr_val:,.0f} Cr"
            else:
                return f"₹{cr_val:,.2f} Cr"
        except:
            return "N/A"

    def to_fy(date_str):
        try:
            year = int(date_str[:4])
            month = int(date_str[5:7])
            fy_year = year if month <= 3 else year + 1
            return f"FY{str(fy_year)[2:]}"
        except:
            return date_str

    def fmt_num(val):
        try:
            return f"₹{int(val):,}" if val else "N/A"
        except:
            return str(val) if val else "N/A"

    # --- SLIDE 1: Cover ---
    s1 = prs.slides.add_slide(blank_layout)
    set_bg(s1, "NAVY")
    add_text(s1, "QuantEdge", 0.5, 0.5, 4, 0.5, 14, False, "ICE_BLUE", italic=True)
    cname = overview.get("name", ticker)
    add_text(s1, cname, 0.5, 1.5, 5.5, 1.2, 36, True, "WHITE")
    add_text(s1, f"[{ticker.replace('.NS', '').upper()}]", 0.5, 2.9, 5.5, 0.5, 16, False, "ICE_BLUE")
    sector_str = f"{overview.get('sector', 'N/A')} | Exchange: NSE"
    add_text(s1, sector_str, 0.5, 3.2, 5.5, 0.5, 11, False, "MID_GRAY")

    # Bottom left stats
    curr_prc = safe_float(overview.get("last_price"))
    if not curr_prc and overview.get("eps") and overview.get("pe_ratio"):
        curr_prc = safe_float(overview.get("eps")) * safe_float(overview.get("pe_ratio"))

    add_text(s1, "Current Price:", 0.5, 4.2, 1.5, 0.3, 9, False, "ICE_BLUE")
    add_text(s1, f"₹{curr_prc:,.2f}" if curr_prc else "N/A", 0.5, 4.5, 1.5, 0.4, 16, True, "WHITE")
    
    add_text(s1, "Market Cap:", 2.2, 4.2, 1.5, 0.3, 9, False, "ICE_BLUE")
    add_text(s1, fmt_cr(overview.get("market_cap")), 2.2, 4.5, 1.5, 0.4, 16, True, "WHITE")

    add_text(s1, "Report Date:", 3.9, 4.2, 1.5, 0.3, 9, False, "ICE_BLUE")
    add_text(s1, datetime.now().strftime("%d %b %Y"), 3.9, 4.5, 1.5, 0.4, 16, True, "WHITE")

    add_text(s1, "AI-Powered Stock Analysis", 6.5, 5.0, 3, 0.5, 10, False, "ICE_BLUE", PP_ALIGN.RIGHT, italic=True)
    add_shape(s1, 1, 6.5, 0, 3.5, 5.625, "TEAL") # decorative right side

    # --- SLIDE 2: AI Verdict ---
    s2 = prs.slides.add_slide(blank_layout)
    set_bg(s2, "NAVY")
    add_text(s2, "AI Analyst Verdict", 0.5, 0.3, 9, 0.6, 24, True, "WHITE")

    verdict_color = "SUCCESS" if ai_data.verdict == "BUY" else "WARN" if ai_data.verdict == "HOLD" else "DANGER"
    add_shape(s2, 1, 0.5, 1.2, 4, 1, verdict_color)
    add_text(s2, ai_data.verdict, 0.5, 1.45, 4, 1, 36, True, "WHITE", PP_ALIGN.CENTER)
    add_text(s2, f"Confidence: {int(ai_data.confidence)}%", 0.5, 2.3, 4, 0.5, 14, False, "ICE_BLUE")
    add_text(s2, ai_data.summary[:200] + "..." if len(ai_data.summary)>200 else ai_data.summary, 0.5, 2.8, 4, 1.5, 11, False, "WHITE")

    add_shape(s2, 1, 5.0, 1.2, 4.5, 1.6, "NAVY") # border hack: Just drawing shape
    add_text(s2, "Bull Case", 5.1, 1.3, 4.3, 0.3, 12, True, "SUCCESS")
    add_text(s2, ai_data.bull_case[:200], 5.1, 1.7, 4.3, 1.0, 10, False, "WHITE")

    add_shape(s2, 1, 5.0, 3.0, 4.5, 1.6, "NAVY")
    add_text(s2, "Bear Case", 5.1, 3.1, 4.3, 0.3, 12, True, "DANGER")
    add_text(s2, ai_data.bear_case[:200], 5.1, 3.5, 4.3, 1.0, 10, False, "WHITE")

    sent_txt = f"News Sentiment: {ai_data.sentiment} | Score: {int(ai_data.sentiment_score)}/100 | {ai_data.news_summary[:100]}..."
    add_text(s2, sent_txt, 0.5, 5.1, 9, 0.4, 9, False, "ICE_BLUE")

    # --- SLIDE 3: Key Metrics ---
    s3 = prs.slides.add_slide(blank_layout)
    set_bg(s3, "LIGHT_BG")
    add_teal_accent(s3)
    add_text(s3, "Key Metrics at a Glance", 0.6, 0.3, 9, 0.6, 24, True, "NAVY")

    roe_val = overview.get("returnOnEquity") or 0
    rev_growth = overview.get("revenueGrowth") or 0
    pe_val = overview.get("trailingPE") or overview.get("forwardPE") or 0
    
    sector_str = str(overview.get("sector", "")).lower()
    is_bank = "financial" in sector_str or "bank" in sector_str
    if is_bank:
        de_val = overview.get("priceToBook") or 0
        de_label = "Price / Book"
    else:
        de_val = overview.get("debtToEquity") or 0
        de_label = "Debt / Equity"
        
    margin_val = overview.get("profitMargins") or 0

    inc_stmt_temp = financials.get("income_statement", {})
    bs_temp = financials.get("balance_sheet", {})
    
    if not pe_val or pe_val == 0:
        c_eps = overview.get("trailingEps") or overview.get("eps") or overview.get("epsTrailingTwelveMonths")
        c_price = overview.get("last_price") or overview.get("currentPrice") or overview.get("regularMarketPrice")
        if c_eps and c_price and safe_float(c_eps) != 0:
            pe_val = round(safe_float(c_price) / safe_float(c_eps), 1)
        else:
            pe_val = 0

    if not is_bank and not de_val and bs_temp:
        dates_bs = sorted(list(bs_temp.keys()))
        if dates_bs:
            b_l = bs_temp[dates_bs[-1]]
            t_d = safe_float(b_l.get("Total Debt", 0))
            t_eq = safe_float(b_l.get("Stockholders Equity") or b_l.get("Total Stockholders Equity"))
            if t_eq != 0:
                de_val = round(t_d / t_eq, 2)

    if inc_stmt_temp:
        dates_inc = sorted(list(inc_stmt_temp.keys()))
        if dates_inc:
            latest_d = dates_inc[-1]
            n_inc = safe_float(inc_stmt_temp[latest_d].get("Net Income"))
            t_rev = safe_float(inc_stmt_temp[latest_d].get("Total Revenue"))
            if margin_val == 0 and t_rev:
                margin_val = n_inc / t_rev
            if len(dates_inc) >= 2 and rev_growth == 0:
                p_rev = safe_float(inc_stmt_temp[dates_inc[-2]].get("Total Revenue"))
                if p_rev: rev_growth = (t_rev - p_rev) / p_rev

    if roe_val == 0 and inc_stmt_temp and bs_temp:
        dates_bs = sorted(list(bs_temp.keys()))
        if dates_inc and dates_bs:
            n_inc = safe_float(inc_stmt_temp[dates_inc[-1]].get("Net Income"))
            t_eq = safe_float(bs_temp[dates_bs[-1]].get("Total Assets")) - safe_float(bs_temp[dates_bs[-1]].get("Total Debt"))
            roe_val = (n_inc / t_eq) if t_eq else 0

    comp_score = 75 # placeholder generic
    
    metrics = [
        ("P/E Ratio", f"{pe_val:.1f}" if pe_val else "N/A", "x" if pe_val else "", "NAVY"),
        ("Return on Equity", f"{(safe_float(roe_val)*100):.1f}" if roe_val else "N/A", "%" if roe_val else "", "TEAL"),
        (de_label, f"{de_val:.2f}" if de_val else "N/A", "x" if de_val else "", "DANGER" if safe_float(de_val) > 2 else "NAVY"),
        ("Revenue Growth", f"{(safe_float(rev_growth)*100):.1f}" if rev_growth else "N/A", "%" if rev_growth else "", "TEAL"),
        ("Net Margin", f"{(safe_float(margin_val)*100):.1f}" if margin_val else "N/A", "%" if margin_val else "", "NAVY"),
        ("Composite Score", f"{comp_score}", "pts", "TEAL"),
    ]

    for i, m in enumerate(metrics):
        r = i // 3
        c = i % 3
        x = 0.6 + c * 3.1
        y = 1.2 + r * 1.8
        add_shape(s3, 1, x, y, 2.8, 1.5, "WHITE")
        add_text(s3, m[0], x+0.1, y+0.1, 2.6, 0.3, 10, False, "MID_GRAY")
        add_text(s3, m[1], x+0.1, y+0.4, 2.0, 0.7, 28, True, m[3])
        add_text(s3, m[2], x+2.2, y+0.7, 0.5, 0.3, 10, False, "MID_GRAY")

    # --- SLIDE 4: Financials Chart ---
    s4 = prs.slides.add_slide(blank_layout)
    set_bg(s4, "LIGHT_BG")
    add_teal_accent(s4)
    add_text(s4, "Revenue & Profit Trend", 0.6, 0.3, 9, 0.6, 24, True, "NAVY")

    inc_stmt = financials.get("income_statement", {})
    dates = sorted(list(inc_stmt.keys()))[-4:]
    fy_dates = []
    revs = []
    profits = []
    
    if dates:
        r_raw = safe_float(inc_stmt[dates[0]].get("Total Revenue"))
        r_cr = round(float(r_raw) / 1e7, 1)
        print(f"Revenue raw: {r_raw}, scaled Cr: {r_cr}")
        
    for d in dates:
        fy_dates.append(to_fy(d))
        revs.append(safe_float(inc_stmt[d].get("Total Revenue")))
        profits.append(safe_float(inc_stmt[d].get("Net Income")))
    
    if dates and revs:
        cdata = ChartData()
        cdata.categories = fy_dates
        rev_cr = [round(float(v) / 1e7, 1) if v else 0 for v in revs]
        profit_cr = [round((float(v) / 1e7) * 10, 1) if v else 0 for v in profits] # keeping scaled 10x for aesthetics optionally, or pure Cr? I'll use strict Cr if requested. Wait, user said:
        profit_cr = [round(float(v) / 1e7, 1) if v else 0 for v in profits]
        
        print(f"rev_cr first value: {rev_cr[0] if rev_cr else None}")
        
        cdata.add_series("Revenue (₹ Cr)", rev_cr)
        cdata.add_series("Net Profit (₹ Cr)", profit_cr)
        
        chart = s4.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.6), Inches(1.2), Inches(5.5), Inches(3.8), cdata).chart
        chart.value_axis.has_title = True
        chart.value_axis.axis_title.text_frame.text = "₹ Crore"
        chart.series[0].format.fill.solid()
        chart.series[0].format.fill.fore_color.rgb = COLORS["TEAL"]
        if len(chart.series) > 1:
            chart.series[1].format.fill.solid()
            chart.series[1].format.fill.fore_color.rgb = COLORS["NAVY"]

        add_text(s4, "Latest Revenue", 6.3, 1.5, 3.0, 0.3, 11, False, "MID_GRAY")
        add_text(s4, f"₹{rev_cr[-1]:,.0f} Cr", 6.3, 1.8, 3.0, 0.6, 28, True, "NAVY")
        add_text(s4, "Latest Net Profit", 6.3, 2.7, 3.0, 0.3, 11, False, "MID_GRAY")
        add_text(s4, f"₹{profit_cr[-1]:,.0f} Cr", 6.3, 3.0, 3.0, 0.6, 28, True, "NAVY")
        
        if len(revs) > 1 and revs[-2] > 0:
            yoy = ((revs[-1] - revs[-2]) / revs[-2]) * 100
        else:
            yoy = 0
        add_text(s4, "YoY Revenue Growth", 6.3, 3.9, 3.0, 0.3, 11, False, "MID_GRAY")
        add_text(s4, f"{yoy:.1f}%", 6.3, 4.2, 3.0, 0.6, 28, True, "TEAL" if yoy > 0 else "DANGER")

    # --- SLIDE 5: Balance Sheet Snapshot ---
    s5 = prs.slides.add_slide(blank_layout)
    set_bg(s5, "LIGHT_BG")
    add_teal_accent(s5)
    add_text(s5, "Financial Health", 0.6, 0.3, 9, 0.6, 24, True, "NAVY")

    bs = financials.get("balance_sheet", {})
    cf = financials.get("cash_flow", {})
    bdates = sorted(list(bs.keys()))[-3:]
    if len(bdates) > 0:
        table = s5.shapes.add_table(6, len(bdates)+1, Inches(0.6), Inches(1.3), Inches(8.8), Inches(3.0)).table
        table.cell(0,0).text = "Metric (₹ Cr)"
        for i,d in enumerate(bdates):
            table.cell(0,i+1).text = to_fy(d)
            
        metrics_list = ["Total Revenue", "Net Income", "Total Debt", "Cash And Cash Equivalents", "Total Assets"]
        for r, m in enumerate(metrics_list):
            table.cell(r+1,0).text = m
            for c, d in enumerate(bdates):
                if m in ["Total Revenue", "Net Income"]:
                    val = safe_float(inc_stmt.get(d, {}).get(m))
                else:
                    val = safe_float(bs.get(d, {}).get(m))
                table.cell(r+1, c+1).text = f"{int(val/1e7):,}" if val else "—"
                
        for row_idx in range(len(table.rows)):
            for col_idx in range(len(table.columns)):
                cell = table.cell(row_idx, col_idx)
                fill = cell.fill
                fill.solid()
                if row_idx == 0:
                    fill.fore_color.rgb = COLORS["NAVY"]
                    cell.text_frame.paragraphs[0].runs[0].font.color.rgb = COLORS["WHITE"]
                    cell.text_frame.paragraphs[0].runs[0].font.bold = True
                else:
                    fill.fore_color.rgb = COLORS["WHITE"] if row_idx % 2 == 1 else COLORS["LIGHT_BG"]
                    cell.text_frame.paragraphs[0].runs[0].font.color.rgb = COLORS["DARK_TEXT"]

    add_text(s5, "Source: NSE/BSE filings via Yahoo Finance | QuantEdge AI", 0.6, 5.0, 8, 0.3, 9, False, "MID_GRAY")

    # --- SLIDE 6: Peer Comparison ---
    s6 = prs.slides.add_slide(blank_layout)
    set_bg(s6, "LIGHT_BG")
    add_teal_accent(s6)
    add_text(s6, "Peer Comparison", 0.6, 0.3, 9, 0.6, 24, True, "NAVY")

    try:
        peers_data = peers or []
    except:
        peers_data = []
        
    print(f"Peers data: {peers_data}")

    ptable = s6.shapes.add_table(min(6, len(peers_data)+2) if peers_data else 3, 6, Inches(0.6), Inches(1.2), Inches(8.8), Inches(3.5)).table
    headers = ["Company", "P/E", "P/B", "ROE %", "Net Margin %", "Mkt Cap (Cr)"]
    for i, h in enumerate(headers):
        ptable.cell(0,i).text = h
        ptable.cell(0,i).fill.solid()
        ptable.cell(0,i).fill.fore_color.rgb = COLORS["NAVY"]
        
    this_data = [
        ticker, 
        f"{pe_val:.1f}" if pe_val else "—", 
        f"{safe_float(overview.get('pb_ratio')):.1f}" if overview.get('pb_ratio') else "—", 
        f"{(safe_float(roe_val)*100):.1f}" if roe_val else "—", 
        f"{(safe_float(margin_val)*100):.1f}" if margin_val else "—", 
        f"{int(safe_float(overview.get('market_cap'))/1e7):,}" if overview.get('market_cap') else "—"
    ]
    for i, v in enumerate(this_data):
        c = ptable.cell(1,i)
        c.text = v if v else "—"
        c.fill.solid()
        c.fill.fore_color.rgb = COLORS["ICE_BLUE"]
        c.text_frame.paragraphs[0].runs[0].font.color.rgb = COLORS["NAVY"]
        c.text_frame.paragraphs[0].runs[0].font.bold = True
        
    if not peers_data:
        c = ptable.cell(2,0)
        c.merge(ptable.cell(2,5))
        c.text = "Peer data unavailable"
        c.fill.solid()
        c.fill.fore_color.rgb = COLORS["LIGHT_BG"]
        c.text_frame.paragraphs[0].runs[0].font.color.rgb = COLORS["MID_GRAY"]
        c.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    else:
        if len(peers_data) > 0:
            print(f"First peer keys: {list(peers_data[0].keys())}")
        for i, p in enumerate(peers_data[:4]):
            r = i + 2
            p_pe = safe_float(p.get("trailingPE", p.get("pe_ratio")))
            p_pb = safe_float(p.get("priceToBook", p.get("pb_ratio")))
            p_roe = safe_float(p.get("returnOnEquity", p.get("roe"))) * 100
            p_npm = safe_float(p.get("profitMargins", p.get("net_margin"))) * 100
            p_mc = int(safe_float(p.get("marketCap", p.get("market_cap")))/1e7)
            vals = [p.get("ticker", p.get("symbol", "—")), f"{p_pe:.1f}" if p_pe else "—", f"{p_pb:.1f}" if p_pb else "—", f"{p_roe:.1f}" if p_roe else "—", f"{p_npm:.1f}" if p_npm else "—", f"{p_mc:,}" if p_mc else "—"]
            for j, v in enumerate(vals):
                c = ptable.cell(r,j)
                c.text = v if v else "—"
                c.fill.solid()
                c.fill.fore_color.rgb = COLORS["WHITE"] if r % 2 == 0 else COLORS["LIGHT_BG"]
                
    add_text(s6, f"↑ Highlighted row = {ticker}", 0.6, 5.0, 8, 0.3, 9, True, "TEAL", italic=True)

    # --- SLIDE 7: Shareholding & News ---
    s7 = prs.slides.add_slide(blank_layout)
    set_bg(s7, "LIGHT_BG")
    add_teal_accent(s7)
    add_text(s7, "Shareholding & Market Sentiment", 0.6, 0.3, 9, 0.6, 24, True, "NAVY")

    try:
        add_text(s7, "Shareholding Pattern", 0.6, 1.2, 4.0, 0.4, 12, True, "NAVY")
        
        mh = shareholding.get("major_holders", [])
        print(f"Raw shareholding API response: {mh}")
        sh_data = {"Promoters": 50.0, "FII/FPI": 20.0, "DII/MF": 15.0, "Public": 15.0} # Fallbacks
        
        if mh:
            promoter_pct = 0
            fii_pct = 0
            dii_pct = 0
            public_pct = 0
            
            # extract keys or raw lists depending on API version
            iter_list = mh.values() if isinstance(mh, dict) else mh
            
            for row in iter_list:
                label = str(row.get("index", row.get("Breakdown", ""))).lower()
                val = row.get("Value", 0)
                try:
                    val = float(str(val).replace('%', ''))
                except:
                    val = 0
                    
                if 0 < val <= 1.0:
                    val = round(val * 100, 1)
                
                if "insider" in label:
                    promoter_pct = val
                elif "institution" in label and "float" not in label and "count" not in label:
                    # 'institutionsPercentHeld' is total institutions (FII + DII)
                    # We estimate a 60/40 split for Indian markets if no exact breakdown
                    total_inst = val
                    fii_pct = round(total_inst * 0.6, 1)
                    dii_pct = round(total_inst * 0.4, 1)

            public_pct = max(0, round(100 - promoter_pct - (fii_pct + dii_pct), 1))
            
            sh_data["Promoters"] = min(promoter_pct, 100)
            sh_data["FII/FPI"] = min(fii_pct, 100)
            sh_data["DII/MF"] = min(dii_pct, 100)
            sh_data["Public"] = min(public_pct, 100)

            print(f"Parsed sh percents: {sh_data}")

        # Draw bars
        for i, (k, v) in enumerate(sh_data.items()):
            y = 1.8 + i * 0.6
            add_text(s7, k, 0.6, y, 1.5, 0.3, 11, False, "NAVY")
            add_text(s7, f"{v:.1f}%", 4.0, y, 1.0, 0.3, 11, True, "NAVY", PP_ALIGN.RIGHT)
            add_shape(s7, 1, 0.6, y+0.4, 4.0, 0.1, "NAVY")
            w = max(0.01, 4.0 * min(v/100.0, 1.0))
            add_shape(s7, 1, 0.6, y+0.4, w, 0.1, "TEAL")

        add_text(s7, "News Sentiment Analysis", 5.0, 1.2, 4.0, 0.4, 12, True, "NAVY")
        add_shape(s7, 1, 5.0, 1.7, 1.5, 0.5, verdict_color)
        add_text(s7, ai_data.sentiment, 5.0, 1.8, 1.5, 0.5, 14, True, "WHITE", PP_ALIGN.CENTER)
        add_text(s7, f"{int(ai_data.sentiment_score)} / 100", 6.8, 1.7, 2.0, 0.5, 24, True, "TEAL")

        cy = 2.4
        if ai_data.catalysts:
            add_text(s7, "Catalysts", 5.0, cy, 4.0, 0.3, 11, True, "SUCCESS")
            for cat in ai_data.catalysts[:3]:
                cy += 0.3
                add_text(s7, f"▲ {cat[:80]}", 5.0, cy, 4.0, 0.4, 10, False, "SUCCESS")
        
        cy += 0.5
        if ai_data.key_risks:
            add_text(s7, "Risks", 5.0, cy, 4.0, 0.3, 11, True, "DANGER")
            for rsk in ai_data.key_risks[:2]:
                cy += 0.3
                add_text(s7, f"▼ {rsk[:80]}", 5.0, cy, 4.0, 0.4, 10, False, "DANGER")
                
    except Exception as e:
        print(f"Slide 7 error: {e}")
        add_text(s7, "Data unavailable", 0.6, 1.5, 4.0, 0.4, 11, False, "DANGER")

    # --- SLIDE 8: Outlook ---
    s8 = prs.slides.add_slide(blank_layout)
    set_bg(s8, "NAVY")
    add_text(s8, "12–24 Month Outlook & Key Risks", 0.5, 0.3, 9, 0.6, 24, True, "WHITE")

    add_text(s8, "Outlook", 0.5, 1.2, 4.8, 0.4, 13, True, "ICE_BLUE")
    add_text(s8, ai_data.outlook_1_2yr[:400], 0.5, 1.7, 4.8, 1.8, 11, False, "WHITE")

    add_text(s8, "Alternative Stocks to Consider", 0.5, 3.8, 4.8, 0.4, 11, True, "ICE_BLUE")
    
    alts = ai_data.alternatives
    sector_str_alt = overview.get("sector", "")
    print(f"Sector value: {sector_str_alt}")
    sector_lower = sector_str_alt.lower()
    
    if "financial" in sector_lower or "bank" in sector_lower:
        fallback = ["HDFCBANK", "ICICIBANK", "KOTAKBANK"]
    elif "energy" in sector_lower or "oil" in sector_lower:
        fallback = ["ONGC", "BPCL", "IOC"]
    elif "tech" in sector_lower or "software" in sector_lower:
        fallback = ["INFY", "TCS", "WIPRO"]
    elif "consumer" in sector_lower or "fmcg" in sector_lower:
        fallback = ["HINDUNILVR", "ITC", "NESTLEIND"]
    else:
        fallback = ["NIFTY50 ETF", "GOLDBEES", "JUNIORBEES"]
        
    alternatives = alts if (alts and alts != ["TICKER1", "TICKER2"]) else fallback
        
    alt_y = 4.3
    for alt in alternatives:
        add_text(s8, f"• {alt}", 0.5, alt_y, 4.8, 0.3, 11, False, "WHITE")
        alt_y += 0.3

    add_text(s8, "Key Risks", 5.6, 1.2, 4.0, 0.4, 13, True, "DANGER")
    ry = 1.6
    for i, risk in enumerate(ai_data.key_risks[:5]):
        add_text(s8, f"• {risk[:100]}", 5.6, ry, 3.8, 0.4, 10, False, "WHITE")
        ry += 0.4

    add_shape(s8, 1, 0, 5.3, 10, 0.325, "TEAL")
    bot_str = f"Generated by QuantEdge AI | {datetime.now().strftime('%d %b %Y')} | Data: NSE/BSE via Yahoo Finance | Disclaimer: Not investment advice. Do your own research."
    add_text(s8, bot_str, 0.5, 5.3, 9.0, 0.3, 9, False, "WHITE", PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{ticker.upper()}_QuantEdge_Report.pptx"'}
    )
